from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Slide content as list of dicts
slides = [
    {"title": "Insurance Policy Voice Agent", "bullets": [
        "Voice-enabled AI assistant for insurance policy queries",
        "Built for hackathon innovation",
        "Combines RAG, voice tech, and no-hallucination AI",
        "Fast, accurate, and user-friendly",
        "Powered by Google ADK, Gemini, FAISS, FastAPI"
    ]},
    {"title": "Problem Statement – Insurance Policy Confusion", "bullets": [
        "Customers struggle to understand complex policy documents",
        "Traditional support is slow and frustrating",
        "Policy details are buried in lengthy PDFs",
        "High risk of misinterpretation and missed benefits",
        "Need for instant, reliable answers"
    ]},
    {"title": "Why Existing Chatbots Fail", "bullets": [
        "Prone to hallucinations and wrong answers",
        "Generic responses, not policy-specific",
        "Poor handling of voice input/output",
        "Lack of document grounding",
        "Low user trust in critical insurance scenarios"
    ]},
    {"title": "Solution Overview – Voice-Enabled RAG Agent", "bullets": [
        "AI agent answers policy questions using real documents",
        "Voice input and output for natural interaction",
        "Retrieval-Augmented Generation (RAG) ensures accuracy",
        "No hallucinations – answers only from indexed policies",
        "FastAPI backend, browser-based frontend"
    ]},
    {"title": "High-Level Architecture", "bullets": [
        "Frontend: Browser UI with Speech-to-Text & Text-to-Speech",
        "Backend: FastAPI orchestrates RAG and voice flows",
        "AI: Google ADK Agents + Gemini LLMs",
        "Vector DB: FAISS for fast, semantic document search",
        "Modular, scalable, and cloud-ready"
    ]},
    {"title": "[Diagram] High-Level Architecture", "diagram": """
┌────────────────────────────────────────────────────────────────────────────┐
│                    Entry Point                                            │
│              RAGOrchestrator (Root Agent)                                 │
│         - Intelligent Request Routing                                     │
│         - Text Mode (_run_async_impl)                                     │
│         - Voice Mode (_run_live_impl)                                     │
└───────────────┬────────────────────────────────────────────────────────────┘
                │
    ┌───────────┼────────────┬────────────┐
    │           │            │            │
    v           v            v
┌──────────┐ ┌───────────┐ ┌────────────┐
│  Policy  │ │  Search   │ │   Voice    │
│ Manager  │ │ Assistant │ │ Assistant  │
│  Agent   │ │  Agent    │ │   Agent    │
└─────┬────┘ └─────┬─────┘ └─────┬──────┘
      v           v              v
┌──────────┐ ┌───────────┐ ┌────────────┐
│ Policy   │ │  Search   │ │   Voice    │
│  Tools   │ │  Tools    │ │   Tools    │
└──────────┘ └───────────┘ └────────────┘
"""},
    {"title": "Detailed RAG Flow", "bullets": [
        "User query received (text or voice)",
        "Query embedded using Gemini/AI model",
        "FAISS vector search retrieves relevant policy chunks",
        "Context passed to LLM for grounded answer",
        "Only document-backed responses returned"
    ]},
    {"title": "[Diagram] System Flow Diagram", "diagram": """
┌────────────────────┐
│   User Input       │
│  (Text or Voice)   │
└─────────┬──────────┘
          v
┌──────────────────────────────────────────────┐
│         RAG Orchestrator                     │
│  1. Extract text from context                │
│  2. Determine route via keywords             │
└─────────┬────────────────────────────────────┘
          v
    ┌────────────┐
    │  Routing   │
    │   Logic    │
    └────┬───────┘
         v
  Keywords found?
     │
 ┌───┴───┐
 │       │
 v       v
Yes     No
 │       │
 v       v
┌────────────┐   ┌───────────────┐
│  Policy    │   │   Search      │
│  Manager   │   │   Assistant   │
│  Agent     │   │   Agent       │
└────┬───────┘   └─────┬─────────┘
     v                 v
┌────────────┐   ┌───────────────┐
│ List Files │   │  Search       │
│   Tool     │   │  Documents    │
└────┬───────┘   └─────┬─────────┘
     v                 v
┌────────────┐   ┌───────────────┐
│  Return    │   │  Return       │
│  Policy    │   │  Answer from  │
│  List      │   │  Indexed Docs │
└────┬───────┘   └─────┬─────────┘
     v                 v
        ┌───────────────┐
        │ Response to   │
        │ User          │
        └───────────────┘
"""},
    {"title": "Voice Flow", "bullets": [
        "User speaks into browser mic",
        "Speech-to-Text (STT) converts to text",
        "Text sent to backend RAG pipeline",
        "Answer generated and sent back",
        "Text-to-Speech (TTS) reads answer aloud"
    ]},
    {"title": "No-Hallucination Strategy", "bullets": [
        "Strict grounding: answers only from indexed documents",
        "RAG pipeline enforces document context",
        "No open-ended LLM responses",
        "Policy snippets shown as evidence",
        "User trust through transparency"
    ]},
    {"title": "Tech Stack", "bullets": [
        "Frontend: HTML, JS, Web Speech API",
        "Backend: FastAPI (Python), Flask (optional)",
        "AI: Google ADK, Gemini models",
        "Vector DB: FAISS",
        "Infra: Local or cloud deployable"
    ]},
    {"title": "Live Demo Flow", "bullets": [
        "Judge asks a real policy question via mic",
        "System transcribes, searches, and answers instantly",
        "Policy snippet shown as source",
        "Voice output reads answer",
        "End-to-end, judge-friendly experience"
    ]},
    {"title": "Scalability & Enterprise Readiness", "bullets": [
        "Modular microservices architecture",
        "Easy to add new policy documents",
        "Scalable vector search (FAISS)",
        "Secure, compliant data handling",
        "Ready for cloud or on-prem deployment"
    ]},
    {"title": "Future Enhancements", "bullets": [
        "Multilingual voice support",
        "Integration with live insurance systems",
        "Advanced analytics and reporting",
        "Mobile app version",
        "Continuous LLM improvement"
    ]},
    {"title": "Hackathon Impact & Closing", "bullets": [
        "Solves real customer pain point",
        "Demonstrates safe, grounded AI",
        "Voice-first, accessible for all users",
        "Ready for enterprise adoption",
        "Thank you! Questions?"
    ]}
]


from pptx.util import Inches

def create_ppt(filename="VoiceAgent_Hackathon_Presentation.pptx"):
    prs = Presentation()
    for idx, slide in enumerate(slides):
        # Title slide
        if idx == 0:
            layout = prs.slide_layouts[0]
            s = prs.slides.add_slide(layout)
            s.shapes.title.text = slide["title"]
            s.placeholders[1].text = slide["bullets"][0]
        # Insert architecture image
        elif slide["title"].lower().find("architecture") != -1 and "image" not in slide:
            layout = prs.slide_layouts[5]  # Title Only
            s = prs.slides.add_slide(layout)
            s.shapes.title.text = slide["title"]
            img_path = "tools/RAGOrchestrator Agent Tool-architrecture.png"
            left = Inches(1)
            top = Inches(1.5)
            width = Inches(7)
            s.shapes.add_picture(img_path, left, top, width=width)
        # Insert flow diagram image
        elif slide["title"].lower().find("flow diagram") != -1 and "image" not in slide:
            layout = prs.slide_layouts[5]  # Title Only
            s = prs.slides.add_slide(layout)
            s.shapes.title.text = slide["title"]
            img_path = "tools/RAGOrchestrator Agent Tool-flow-diagram.png"
            left = Inches(1)
            top = Inches(1.5)
            width = Inches(7)
            s.shapes.add_picture(img_path, left, top, width=width)
        # Technical/Business explanation slides
        elif slide["title"] == "Tech Stack":
            layout = prs.slide_layouts[1]
            s = prs.slides.add_slide(layout)
            s.shapes.title.text = "Tech Stack: Why These Choices?"
            body = s.placeholders[1].text_frame
            body.text = "• FastAPI: High-performance, async Python web framework.\n  - Technical: Enables fast, scalable APIs for real-time RAG and voice.\n  - Business: Reduces latency, improves user experience.\n\n"
            body.add_paragraph().text = "• FAISS: Facebook AI Similarity Search.\n  - Technical: Efficient vector search for large policy docs.\n  - Business: Enables instant, accurate retrieval, scales to enterprise data.\n\n"
            body.add_paragraph().text = "• Google ADK & Gemini: Advanced LLM and agent orchestration.\n  - Technical: State-of-the-art language understanding and agent routing.\n  - Business: Reduces hallucination, increases trust, future-proofs solution.\n\n"
            body.add_paragraph().text = "• python-pptx: Automated PPT generation.\n  - Technical: Enables hackathon-speed, repeatable slide creation.\n  - Business: Saves time, ensures consistency for demos.\n\n"
            body.add_paragraph().text = "• Web Speech API: Browser-based STT/TTS.\n  - Technical: Seamless voice input/output in any browser.\n  - Business: No install, accessible to all users."
        elif slide["title"] == "Tech Stack":
            continue
        # Default: bullet slides
        elif "bullets" in slide:
            layout = prs.slide_layouts[1]
            s = prs.slides.add_slide(layout)
            s.shapes.title.text = slide["title"]
            body = s.placeholders[1].text_frame
            for bullet in slide["bullets"]:
                p = body.add_paragraph()
                p.text = bullet
                p.level = 0
                p.font.size = Pt(24)
        # Diagrams as text (fallback)
        elif "diagram" in slide:
            layout = prs.slide_layouts[1]
            s = prs.slides.add_slide(layout)
            s.shapes.title.text = slide["title"]
            body = s.placeholders[1].text_frame
            p = body.paragraphs[0]
            p.text = slide["diagram"].strip()
            p.font.size = Pt(12)
            p.font.name = "Courier New"
            p.alignment = PP_ALIGN.LEFT
    # Add a business value summary slide
    layout = prs.slide_layouts[1]
    s = prs.slides.add_slide(layout)
    s.shapes.title.text = "Business Value & Impact"
    body = s.placeholders[1].text_frame
    body.text = (
        "• Reduces customer confusion and support costs\n"
        "• Increases policyholder satisfaction and trust\n"
        "• Scalable to any insurance product or region\n"
        "• Voice-first: accessible to all, including non-technical users\n"
        "• Ready for enterprise deployment and future enhancements"
    )
    prs.save(filename)
    print(f"Presentation saved as {filename}")

if __name__ == "__main__":
    create_ppt()
